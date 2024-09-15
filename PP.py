from pptx import Presentation
from pptx.util import Inches

# Create a new PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Ứng dụng hệ thống quản lý doanh nghiệp ứng dụng ASP.NET"
subtitle.text = "Khóa luận tốt nghiệp\nĐào Quang Việt\nTrường Đại học Kinh Bắc\n2024"

# Slide 2: Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "Bối cảnh: Trong bối cảnh kinh tế hiện nay, các doanh nghiệp ngày càng chú trọng đến việc ứng dụng công nghệ "
    "thông tin vào quản lý và vận hành.\n\n"
    "Mục tiêu: Nghiên cứu và phát triển một giải pháp quản lý toàn diện cho doanh nghiệp dựa trên nền tảng ASP.NET.\n\n"
    "Tầm quan trọng: Giúp doanh nghiệp quản lý các hoạt động hàng ngày hiệu quả, cung cấp các công cụ hỗ trợ ra "
    "quyết định, nâng cao năng suất và khả năng cạnh tranh."
)

# Slide 3: Objective of the Study
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objective of the Study"
content.text = (
    "Mục đích: Phát triển giải pháp quản lý toàn diện dựa trên ASP.NET.\n\n"
    "Mục tiêu cụ thể: Quản lý hiệu quả các hoạt động hàng ngày, cung cấp công cụ hỗ trợ quyết định, nâng cao năng suất."
)
# Continue adding slides to the existing PowerPoint presentation

# Slide 4: ASP.NET Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "ASP.NET Overview"
content.text = (
    "Định nghĩa: ASP.NET là một nền tảng phát triển web từ Microsoft.\n\n"
    "Các tính năng chính: Hỗ trợ MVC, Web Forms, Web API, và SignalR.\n\n"
    "Ưu điểm: Bảo mật, hiệu năng cao, dễ bảo trì, và hỗ trợ mạnh mẽ từ cộng đồng."
)

# Slide 5: Theoretical Background
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Theoretical Background"
content.text = (
    "Khái niệm hệ thống quản lý doanh nghiệp trực tuyến: Hệ thống giúp quản lý các quy trình kinh doanh thông qua giao diện web.\n\n"
    "Giới thiệu mô hình MVC: Mô hình Model-View-Controller giúp tách biệt dữ liệu, giao diện, và logic điều khiển.\n\n"
    "Lợi ích của kiến trúc MVC: Tăng tính linh hoạt, dễ bảo trì, và mở rộng."
)

# Slide 6: ASP.NET MVC
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "ASP.NET MVC"
content.text = (
    "Giải thích ASP.NET MVC: Mô hình phát triển web sử dụng ASP.NET.\n\n"
    "Các thành phần chính: Model, View, Controller.\n\n"
    "Ưu điểm so với ASP.NET truyền thống: Tách biệt rõ ràng các thành phần, dễ kiểm thử và bảo trì."
)

# Slide 7: SQL Server Overview
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "SQL Server Overview"
content.text = (
    "Giới thiệu SQL Server: Hệ quản trị cơ sở dữ liệu từ Microsoft.\n\n"
    "Vai trò của SQL Server trong ứng dụng ASP.NET: Lưu trữ và quản lý dữ liệu cho ứng dụng."
)

# Slide 8: System Design
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "System Design"
content.text = (
    "Tổng quan kiến trúc hệ thống: Mô hình client-server.\n\n"
    "Thiết kế cơ sở dữ liệu: Các bảng và quan hệ.\n\n"
    "Các module chính và chức năng: Quản lý người dùng, quản lý sản phẩm, báo cáo."
)

# Slide 9: Implementation Process
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Implementation Process"
content.text = (
    "Các bước từ phân tích yêu cầu đến triển khai: Thu thập yêu cầu, thiết kế hệ thống, phát triển, kiểm thử, triển khai.\n\n"
    "Các công cụ và công nghệ sử dụng: Visual Studio, SQL Server Management Studio, Git."
)

# Slide 10: System Features
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "System Features"
content.text = (
    "Quản lý người dùng: Đăng ký, đăng nhập, phân quyền.\n\n"
    "Kiểm soát truy cập: Quản lý quyền hạn và truy cập.\n\n"
    "Ghi nhật ký hoạt động: Theo dõi và ghi lại các hoạt động của người dùng."
)

# Slide 11: User Interface
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "User Interface"
content.text = (
    "Thiết kế giao diện thân thiện người dùng: Tối ưu trải nghiệm người dùng.\n\n"
    "Các màn hình chính và quy trình làm việc: Dashboard, quản lý tài khoản, quản lý sản phẩm."
)

# Slide 12: Testing and Evaluation
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Testing and Evaluation"
content.text = (
    "Phương pháp kiểm thử: Kiểm thử đơn vị, kiểm thử tích hợp, kiểm thử hệ thống.\n\n"
    "Kết quả đánh giá hiệu năng: Đáp ứng nhanh, xử lý dữ liệu lớn.\n\n"
    "Phản hồi từ người dùng: Phản hồi tích cực, đề xuất cải tiến."
)

# Slide 13: Case Study
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Case Study"
content.text = (
    "Ví dụ ứng dụng thực tế: Một doanh nghiệp áp dụng hệ thống.\n\n"
    "Lợi ích quan sát được: Tăng năng suất, quản lý hiệu quả.\n\n"
    "Thách thức gặp phải: Vấn đề kỹ thuật, sự chấp nhận của người dùng."
)

# Slide 14: Benefits of the System
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Benefits of the System"
content.text = (
    "Tăng năng suất: Tự động hóa các quy trình.\n\n"
    "Cải thiện quyết định: Cung cấp báo cáo chi tiết.\n\n"
    "Tính linh hoạt và mở rộng: Dễ dàng mở rộng theo nhu cầu doanh nghiệp."
)

# Slide 15: Challenges and Limitations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Challenges and Limitations"
content.text = (
    "Thách thức kỹ thuật: Vấn đề tích hợp, hiệu năng.\n\n"
    "Vấn đề chấp nhận của người dùng: Đào tạo, thay đổi thói quen làm việc.\n\n"
    "Cải tiến tương lai: Đề xuất cải tiến để khắc phục hạn chế hiện tại."
)

# Slide 16: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "Tóm tắt kết quả nghiên cứu: Hệ thống quản lý doanh nghiệp hiệu quả.\n\n"
    "Tác động tổng thể của hệ thống: Nâng cao năng suất và khả năng cạnh tranh của doanh nghiệp."
)

# Slide 17: Recommendations
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Recommendations"
content.text = (
    "Đề xuất cho nghiên cứu tương lai: Nghiên cứu sâu hơn về bảo mật, tích hợp các công nghệ mới.\n\n"
    "Cải tiến tiềm năng cho hệ thống: Ứng dụng AI, cải thiện giao diện người dùng."
)

# Slide 18: Acknowledgements
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Acknowledgements"
content.text = (
    "Cảm ơn người hướng dẫn: TS. Trần Đức Nghĩa.\n\n"
    "Cảm ơn trường đại học: Trường Đại học Kinh Bắc.\n\n"
    "Lời cảm ơn cá nhân: Gia đình, bạn bè, đồng nghiệp."
)

# Slide 19: References
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "References"
content.text = (
    "Danh sách tài liệu tham khảo chính: Liệt kê các tài liệu đã sử dụng trong khóa luận."
)

# Slide 20: Q&A
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Q&A"
content.text = (
    "Mời câu hỏi và thảo luận: Không gian cho khán giả đặt câu hỏi và thảo luận."
)

# Save the presentation
pptx_path = "C:/Users/Admin/Downloads/ASP_NET_Enterprise_Management_System_Presentation_v2.pptx"
prs.save(pptx_path)

pptx_path
